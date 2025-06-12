"""
REPORT GENERATOR - Telefónica del Perú (ESTRUCTURA CORPORATIVA)
===============================================================

Generador de reportes con estructura jerárquica específica de Telefónica:
CARTERA > SERVICIO > VENCIMIENTO

Basado en la estructura de presentaciones corporativas existentes.
"""

import pandas as pd
import numpy as np
from datetime import datetime, date, timedelta
import tempfile
import os
from typing import Dict, List, Optional, Tuple
import logging

# Excel libraries
import openpyxl
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.worksheet.table import Table, TableStyleInfo

# PowerPoint libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR

logger = logging.getLogger(__name__)

class TelefonicaCorporateReportGenerator:
    """
    Generador de reportes corporativos para Telefónica del Perú
    Estructura jerárquica: CARTERA > SERVICIO > VENCIMIENTO
    """
    
    # Colores corporativos Telefónica/Movistar
    COLORS = {
        'movistar_blue': '0066CC',
        'movistar_green': '00A651', 
        'movistar_dark_blue': '003366',
        'telefonica_blue': '0019A5',
        'telefonica_light_blue': '5BB4E5',
        'telefonica_orange': 'FF6600',
        'white': 'FFFFFF',
        'light_gray': 'F2F2F2',
        'dark_gray': '666666'
    }
    
    # Jerarquía de datos
    HIERARCHY = {
        'CARTERA': ['Altas_Nuevas', 'Temprana', 'Fraccionamiento'],
        'SERVICIO': ['FIJA', 'MOVIL'],
        'VENCIMIENTO': ['05', '09', '13', '17', '21', '25', '01']  # Días del mes
    }
    
    # KPIs corporativos específicos
    KPIS = {
        'contactabilidad': '%CONT',
        'contacto_directo': 'CD%', 
        'contacto_indirecto': 'CI%',
        'conversion': '%CONV',
        'tasa_contacto': 'TC%',
        'intensidad': 'INTENSIDAD',
        'tasa_cierre': 'TASA_CIERRE%',
        'convertibilidad': 'CONVERTIBILIDAD%',
        'efectividad': 'EFECTIVIDAD%'
    }
    
    def __init__(self, fecha_inicio: str, fecha_fin: str, mes_actual: str = None, mes_anterior: str = None):
        """
        Inicializar generador corporativo
        
        Args:
            fecha_inicio: Fecha inicio período (YYYY-MM-DD)
            fecha_fin: Fecha fin período (YYYY-MM-DD)
            mes_actual: Nombre del mes actual (ej: "Abril")
            mes_anterior: Nombre del mes anterior (ej: "Marzo")
        """
        self.fecha_inicio = fecha_inicio
        self.fecha_fin = fecha_fin
        self.mes_actual = mes_actual or datetime.strptime(fecha_fin, '%Y-%m-%d').strftime('%B')
        self.mes_anterior = mes_anterior or (datetime.strptime(fecha_fin, '%Y-%m-%d') - timedelta(days=30)).strftime('%B')
        self.fecha_generacion = datetime.now()
        
        # Estructura de datos corporativa
        self.data = {
            'asignacion_cartera': {},      # Slide 2: Asignación por cartera
            'kpis_integrales': {},         # Slide 3: KPIs comparativo
            'kpis_evolucion': {},          # Slide 4: Evolución temporal
            'kpis_contactabilidad': {},    # Slide 5: Análisis contactabilidad
            'kpis_tipos_contacto': {},     # Slide 6: CD vs CI
            'kpis_resultados': {},         # Slide 7: Tasa cierre y conversión
            'kpis_esfuerzo': {},           # Slide 8: Intensidad y efectividad
            'cumplimiento_objetivo': {},   # Slide 9: Resultados vs metas
            'ranking_agentes': {},         # Slide 10: Performance individual
            'estrategia_gestion': {},      # Slide 11: Metodología
            'implementacion_bot': {}       # Slide 12: Tecnología
        }
    
    def load_data_from_processing(self, 
                                gestiones_df: pd.DataFrame,
                                calendario_df: pd.DataFrame,
                                asignacion_df: pd.DataFrame,
                                pagos_df: pd.DataFrame,
                                kpis_campania: List[Dict]) -> None:
        """
        Cargar y procesar datos según estructura corporativa
        """
        logger.info("Procesando datos con estructura corporativa Telefónica")
        
        # Procesar asignación por cartera
        self._process_asignacion_cartera(calendario_df, asignacion_df)
        
        # Procesar KPIs integrales
        self._process_kpis_integrales(gestiones_df, kpis_campania)
        
        # Procesar evolución temporal
        self._process_kpis_evolucion(gestiones_df)
        
        # Procesar análisis de contactabilidad
        self._process_kpis_contactabilidad(gestiones_df)
        
        # Procesar tipos de contacto
        self._process_kpis_tipos_contacto(gestiones_df)
        
        # Procesar resultados
        self._process_kpis_resultados(gestiones_df)
        
        # Procesar esfuerzo y efectividad
        self._process_kpis_esfuerzo(gestiones_df)
        
        # Procesar cumplimiento de objetivos
        self._process_cumplimiento_objetivo(gestiones_df, pagos_df)
        
        # Procesar ranking de agentes
        self._process_ranking_agentes(gestiones_df)
        
        logger.info("Datos corporativos procesados exitosamente")
    
    def _process_asignacion_cartera(self, calendario_df: pd.DataFrame, asignacion_df: pd.DataFrame) -> None:
        """Procesar datos de asignación por cartera (Slide 2)"""
        if calendario_df.empty or asignacion_df.empty:
            return
        
        # Agrupar por cartera y servicio
        asignacion_summary = {}
        
        for _, row in asignacion_df.iterrows():
            cartera = row.get('tipo_cartera', 'Otro')
            servicio = row.get('servicio_normalizado', 'FIJA')
            clientes = row.get('clientes_asignados', 0)
            
            if cartera not in asignacion_summary:
                asignacion_summary[cartera] = {'FIJA': 0, 'MOVIL': 0, 'total': 0}
            
            asignacion_summary[cartera][servicio] += clientes
            asignacion_summary[cartera]['total'] += clientes
        
        # Procesar por vencimientos (extraer del archivo)
        vencimientos_data = {}
        for _, row in calendario_df.iterrows():
            archivo = row.get('archivo', '')
            vencimiento = str(row.get('vencimiento', 0)).zfill(2)
            suma_lineas = row.get('suma_lineas', 0)
            cartera = row.get('tipo_cartera', 'Otro')
            
            if vencimiento not in vencimientos_data:
                vencimientos_data[vencimiento] = {}
            
            if cartera not in vencimientos_data[vencimiento]:
                vencimientos_data[vencimiento][cartera] = 0
            
            vencimientos_data[vencimiento][cartera] += suma_lineas
        
        self.data['asignacion_cartera'] = {
            'resumen_cartera': asignacion_summary,
            'por_vencimiento': vencimientos_data,
            'comparativa_mensual': {
                'mes_actual': self.mes_actual,
                'mes_anterior': self.mes_anterior
            }
        }
    
    def _process_kpis_integrales(self, gestiones_df: pd.DataFrame, kpis_campania: List[Dict]) -> None:
        """Procesar KPIs integrales comparativos (Slide 3)"""
        if gestiones_df.empty:
            return
        
        # Calcular KPIs por cartera
        kpis_por_cartera = {}
        
        for cartera in self.HIERARCHY['CARTERA']:
            gestiones_cartera = gestiones_df[gestiones_df['tipo_cartera'] == cartera]
            
            if not gestiones_cartera.empty:
                total_gestiones = len(gestiones_cartera)
                contactos_efectivos = len(gestiones_cartera[gestiones_cartera['contactabilidad'] == 'CONTACTO_EFECTIVO'])
                contactos_directos = contactos_efectivos  # Simplificado por ahora
                contactos_indirectos = len(gestiones_cartera[gestiones_cartera['contactabilidad'] == 'CONTACTO_NO_EFECTIVO'])
                compromisos = len(gestiones_cartera[gestiones_cartera['es_pdp'] == 'SI'])
                
                kpis_por_cartera[cartera] = {
                    '%CONT': round(contactos_efectivos / total_gestiones * 100, 2) if total_gestiones > 0 else 0,
                    'CD%': round(contactos_directos / total_gestiones * 100, 2) if total_gestiones > 0 else 0,
                    'CI%': round(contactos_indirectos / total_gestiones * 100, 2) if total_gestiones > 0 else 0,
                    '%CONV': round(compromisos / contactos_efectivos * 100, 2) if contactos_efectivos > 0 else 0,
                    'INTENSIDAD': round(total_gestiones / gestiones_cartera['cod_luna'].nunique(), 2) if gestiones_cartera['cod_luna'].nunique() > 0 else 0,
                    'EFECTIVIDAD%': round(gestiones_cartera['monto_compromiso'].sum() / gestiones_cartera['monto_exigible'].sum() * 100, 2) if gestiones_cartera['monto_exigible'].sum() > 0 else 0
                }
        
        self.data['kpis_integrales'] = kpis_por_cartera
    
    def _process_kpis_evolucion(self, gestiones_df: pd.DataFrame) -> None:
        """Procesar evolución temporal de KPIs (Slide 4)"""
        if gestiones_df.empty:
            return
        
        # Agrupar por fecha y calcular KPIs diarios
        gestiones_df['fecha'] = pd.to_datetime(gestiones_df['date']).dt.date
        evolucion_diaria = {}
        
        for fecha in gestiones_df['fecha'].unique():
            gestiones_dia = gestiones_df[gestiones_df['fecha'] == fecha]
            
            total_gestiones = len(gestiones_dia)
            contactos_efectivos = len(gestiones_dia[gestiones_dia['contactabilidad'] == 'CONTACTO_EFECTIVO'])
            
            evolucion_diaria[fecha.strftime('%d')] = {
                'CONTACTABILIDAD_%': round(contactos_efectivos / total_gestiones * 100, 2) if total_gestiones > 0 else 0,
                'CD_%': round(contactos_efectivos / total_gestiones * 100, 2) if total_gestiones > 0 else 0,  # Simplificado
                'CI_%': round(len(gestiones_dia[gestiones_dia['contactabilidad'] == 'CONTACTO_NO_EFECTIVO']) / total_gestiones * 100, 2) if total_gestiones > 0 else 0,
                'TC_%': round(contactos_efectivos / total_gestiones * 100, 2) if total_gestiones > 0 else 0,  # Simplificado
                'CONVERSION_%': round(len(gestiones_dia[gestiones_dia['es_pdp'] == 'SI']) / contactos_efectivos * 100, 2) if contactos_efectivos > 0 else 0,
                'INTENSIDAD': round(total_gestiones / gestiones_dia['cod_luna'].nunique(), 2) if gestiones_dia['cod_luna'].nunique() > 0 else 0
            }
        
        self.data['kpis_evolucion'] = evolucion_diaria
    
    def _process_kpis_contactabilidad(self, gestiones_df: pd.DataFrame) -> None:
        """Procesar análisis de contactabilidad por servicio (Slide 5)"""
        contactabilidad_por_servicio = {}
        
        for servicio in self.HIERARCHY['SERVICIO']:
            gestiones_servicio = gestiones_df[gestiones_df['servicio'] == servicio]
            
            if not gestiones_servicio.empty:
                # Por vencimiento
                contactabilidad_vcto = {}
                for vcto in ['05', '09', '13', '17']:
                    # Filtrar por día de vencimiento (simplificado)
                    gestiones_vcto = gestiones_servicio[gestiones_servicio['dias_desde_asignacion'] <= int(vcto)]
                    
                    if not gestiones_vcto.empty:
                        total = len(gestiones_vcto)
                        contactos = len(gestiones_vcto[gestiones_vcto['contactabilidad'] == 'CONTACTO_EFECTIVO'])
                        contactabilidad_vcto[f'VCTO_{vcto}'] = round(contactos / total * 100, 2) if total > 0 else 0
                
                contactabilidad_por_servicio[servicio] = contactabilidad_vcto
        
        self.data['kpis_contactabilidad'] = contactabilidad_por_servicio
    
    def _process_kpis_tipos_contacto(self, gestiones_df: pd.DataFrame) -> None:
        """Procesar tipos de contacto CD vs CI (Slide 6)"""
        tipos_contacto = {}
        
        for servicio in self.HIERARCHY['SERVICIO']:
            gestiones_servicio = gestiones_df[gestiones_df['servicio'] == servicio]
            
            if not gestiones_servicio.empty:
                total = len(gestiones_servicio)
                cd = len(gestiones_servicio[gestiones_servicio['contactabilidad'] == 'CONTACTO_EFECTIVO'])
                ci = len(gestiones_servicio[gestiones_servicio['contactabilidad'] == 'CONTACTO_NO_EFECTIVO'])
                
                tipos_contacto[servicio] = {
                    'CONTACTO_DIRECTO_%': round(cd / total * 100, 2) if total > 0 else 0,
                    'CONTACTO_INDIRECTO_%': round(ci / total * 100, 2) if total > 0 else 0
                }
        
        self.data['kpis_tipos_contacto'] = tipos_contacto
    
    def _process_kpis_resultados(self, gestiones_df: pd.DataFrame) -> None:
        """Procesar KPIs de resultados (Slide 7)"""
        resultados = {}
        
        for servicio in self.HIERARCHY['SERVICIO']:
            gestiones_servicio = gestiones_df[gestiones_df['servicio'] == servicio]
            
            if not gestiones_servicio.empty:
                contactos_efectivos = len(gestiones_servicio[gestiones_servicio['contactabilidad'] == 'CONTACTO_EFECTIVO'])
                compromisos = len(gestiones_servicio[gestiones_servicio['es_pdp'] == 'SI'])
                
                resultados[servicio] = {
                    'TASA_CIERRE_%': round(compromisos / contactos_efectivos * 100, 2) if contactos_efectivos > 0 else 0,
                    'CONVERTIBILIDAD_%': round(compromisos / len(gestiones_servicio) * 100, 2) if len(gestiones_servicio) > 0 else 0
                }
        
        self.data['kpis_resultados'] = resultados
    
    def _process_kpis_esfuerzo(self, gestiones_df: pd.DataFrame) -> None:
        """Procesar KPIs de esfuerzo y eficiencia (Slide 8)"""
        esfuerzo = {}
        
        for servicio in self.HIERARCHY['SERVICIO']:
            gestiones_servicio = gestiones_df[gestiones_df['servicio'] == servicio]
            
            if not gestiones_servicio.empty:
                clientes_unicos = gestiones_servicio['cod_luna'].nunique()
                total_gestiones = len(gestiones_servicio)
                monto_recuperado = gestiones_servicio['monto_compromiso'].sum()
                monto_exigible = gestiones_servicio['monto_exigible'].sum()
                
                esfuerzo[servicio] = {
                    'INTENSIDAD': round(total_gestiones / clientes_unicos, 2) if clientes_unicos > 0 else 0,
                    'EFECTIVIDAD_%': round(monto_recuperado / monto_exigible * 100, 2) if monto_exigible > 0 else 0
                }
        
        self.data['kpis_esfuerzo'] = esfuerzo
    
    def _process_cumplimiento_objetivo(self, gestiones_df: pd.DataFrame, pagos_df: pd.DataFrame) -> None:
        """Procesar cumplimiento de objetivos (Slide 9)"""
        # Objetivos simulados (se pueden configurar externamente)
        objetivos_base = {
            'FIJA': {'05': 15.0, '09': 18.0, '13': 20.0, '17': 22.0},
            'MOVIL': {'05': 12.0, '09': 15.0, '13': 17.0, '17': 19.0}
        }
        
        cumplimiento = {}
        cumplimiento_general = 0
        peso_total = 0
        
        for servicio in self.HIERARCHY['SERVICIO']:
            cumplimiento[servicio] = {}
            
            for vcto in ['05', '09', '13', '17']:
                # Filtrar gestiones por servicio y vencimiento
                gestiones_filtradas = gestiones_df[
                    (gestiones_df['servicio'] == servicio) & 
                    (gestiones_df['dias_desde_asignacion'] <= int(vcto))
                ]
                
                if not gestiones_filtradas.empty:
                    monto_recuperado = gestiones_filtradas['monto_compromiso'].sum()
                    monto_exigible = gestiones_filtradas['monto_exigible'].sum()
                    recupero_real = round(monto_recuperado / monto_exigible * 100, 2) if monto_exigible > 0 else 0
                    
                    objetivo = objetivos_base.get(servicio, {}).get(vcto, 15.0)
                    cumplimiento_pct = round(recupero_real / objetivo * 100, 2) if objetivo > 0 else 0
                    peso = monto_exigible  # Peso basado en monto exigible
                    
                    cumplimiento[servicio][f'VCTO_{vcto}'] = {
                        'recupero_esperado': objetivo,
                        'recupero_real': recupero_real,
                        'cumplimiento_%': cumplimiento_pct,
                        'peso': peso
                    }
                    
                    # Acumular para cumplimiento general
                    cumplimiento_general += cumplimiento_pct * peso
                    peso_total += peso
        
        # Calcular cumplimiento general
        cumplimiento_general_pct = round(cumplimiento_general / peso_total, 2) if peso_total > 0 else 0
        
        self.data['cumplimiento_objetivo'] = {
            'por_servicio': cumplimiento,
            'cumplimiento_general': cumplimiento_general_pct
        }
    
    def _process_ranking_agentes(self, gestiones_df: pd.DataFrame) -> None:
        """Procesar ranking de agentes (Slide 10)"""
        if gestiones_df.empty:
            return
        
        # Agrupar por agente
        ranking = []
        agentes_stats = gestiones_df.groupby('ejecutivo_homologado').agg({
            'cod_luna': 'count',  # Total gestiones
            'contactabilidad': lambda x: (x == 'CONTACTO_EFECTIVO').sum(),  # CEF
            'contactabilidad': lambda x: (x == 'CONTACTO_NO_EFECTIVO').sum(),  # NEF
            'es_pdp': lambda x: (x == 'SI').sum(),  # Compromisos
            'monto_compromiso': 'sum'  # Monto pagado
        }).round(2)
        
        # Procesar cada agente
        for agente, stats in agentes_stats.iterrows():
            if agente and agente != 'AGENTE NO IDENTIFICADO':
                gestiones_agente = gestiones_df[gestiones_df['ejecutivo_homologado'] == agente]
                
                total_gestiones = len(gestiones_agente)
                cef = len(gestiones_agente[gestiones_agente['contactabilidad'] == 'CONTACTO_EFECTIVO'])
                nef = len(gestiones_agente[gestiones_agente['contactabilidad'] == 'CONTACTO_NO_EFECTIVO'])
                compromisos = len(gestiones_agente[gestiones_agente['es_pdp'] == 'SI'])
                monto = gestiones_agente['monto_compromiso'].sum()
                
                convertibilidad = round(compromisos / total_gestiones * 100, 2) if total_gestiones > 0 else 0
                tasa_cierre = round(compromisos / cef * 100, 2) if cef > 0 else 0
                
                ranking.append({
                    'agente': agente,
                    'gestiones': total_gestiones,
                    'cef': cef,
                    'nef': nef,
                    'monto_pagado': round(monto, 2),
                    'convertibilidad_%': convertibilidad,
                    'tasa_cierre_%': tasa_cierre
                })
        
        # Ordenar por convertibilidad
        ranking.sort(key=lambda x: x['convertibilidad_%'], reverse=True)
        
        # Asignar ranking y cuartiles
        for i, agente_data in enumerate(ranking[:20]):  # Top 20
            agente_data['ranking'] = i + 1
            
            # Asignar cuartil
            if i < 5:
                agente_data['cuartil'] = 'Q1 - Excelente'
            elif i < 10:
                agente_data['cuartil'] = 'Q2 - Bueno'
            elif i < 15:
                agente_data['cuartil'] = 'Q3 - Regular'
            else:
                agente_data['cuartil'] = 'Q4 - Necesita Mejora'
        
        self.data['ranking_agentes'] = ranking[:20]
    
    def generate_powerpoint_corporate(self, output_path: str = None) -> str:
        """
        Generar presentación corporativa con los 13 slides específicos
        """
        if output_path is None:
            timestamp = self.fecha_generacion.strftime('%Y%m%d_%H%M%S')
            output_path = f"Informe_Gestion_Cobranza_Telefonica_{timestamp}.pptx"
        
        logger.info(f"Generando presentación corporativa: {output_path}")
        
        # Crear presentación
        prs = Presentation()
        
        # Generar todos los slides corporativos
        self._create_slide_01_portada(prs)
        self._create_slide_02_asignacion_temprana(prs)
        self._create_slide_03_kpis_integrales(prs)
        self._create_slide_04_kpis_evolucion(prs)
        self._create_slide_05_kpis_contactabilidad(prs)
        self._create_slide_06_kpis_tipos_contacto(prs)
        self._create_slide_07_kpis_resultados(prs)
        self._create_slide_08_kpis_esfuerzo(prs)
        self._create_slide_09_cumplimiento_objetivo(prs)
        self._create_slide_10_ranking_agentes(prs)
        self._create_slide_11_estrategia_gestion(prs)
        self._create_slide_12_implementacion_bot(prs)
        self._create_slide_13_cierre(prs)
        
        # Guardar presentación
        prs.save(output_path)
        logger.info(f"Presentación corporativa generada: {output_path}")
        
        return output_path
    
    def _create_slide_01_portada(self, prs: Presentation) -> None:
        """Slide 1: Portada corporativa"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Fondo azul corporativo
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(self.COLORS['movistar_blue'])
        
        # Título principal
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.text = "INFORME GESTIÓN COBRANZA – TELEFÓNICA PERÚ"
        
        # Formatear título
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(self.COLORS['white'])
        p.alignment = PP_ALIGN.CENTER
        
        # Subtítulo con período
        left = Inches(1)
        top = Inches(4.5)
        width = Inches(8)
        height = Inches(1)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        tf_sub = subtitle_box.text_frame
        tf_sub.text = f"Período: {self.fecha_inicio} - {self.fecha_fin}"
        
        p_sub = tf_sub.paragraphs[0]
        p_sub.font.size = Pt(18)
        p_sub.font.color.rgb = RGBColor.from_string(self.COLORS['white'])
        p_sub.alignment = PP_ALIGN.CENTER
        
        # Logo simulado (texto)
        left = Inches(4)
        top = Inches(6)
        width = Inches(2)
        height = Inches(1)
        logo_box = slide.shapes.add_textbox(left, top, width, height)
        tf_logo = logo_box.text_frame
        tf_logo.text = "M"
        
        p_logo = tf_logo.paragraphs[0]
        p_logo.font.size = Pt(72)
        p_logo.font.bold = True
        p_logo.font.color.rgb = RGBColor.from_string(self.COLORS['movistar_green'])
        p_logo.alignment = PP_ALIGN.CENTER
    
    def _create_slide_02_asignacion_temprana(self, prs: Presentation) -> None:
        """Slide 2: Asignación Temprana - Comparativa de volumen y valor"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "ASIGNACIÓN TEMPRANA"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Resumen de asignación
        asignacion = self.data.get('asignacion_cartera', {})
        resumen = asignacion.get('resumen_cartera', {})
        
        tf.text = f"COMPARATIVA {self.mes_anterior.upper()} vs {self.mes_actual.upper()}"
        
        # Datos por cartera
        for cartera, datos in resumen.items():
            p = tf.add_paragraph()
            p.text = f"• {cartera}: {datos['total']:,} clientes (FIJA: {datos['FIJA']:,}, MÓVIL: {datos['MOVIL']:,})"
            p.font.size = Pt(14)
        
        # Información de vencimientos
        vencimientos = asignacion.get('por_vencimiento', {})
        if vencimientos:
            p = tf.add_paragraph()
            p.text = f"\nDistribución por Vencimiento:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            for vcto, datos in vencimientos.items():
                total_vcto = sum(datos.values())
                p = tf.add_paragraph()
                p.text = f"• VCTO {vcto}: {total_vcto:,} clientes"
                p.font.size = Pt(12)
    
    def _create_slide_03_kpis_integrales(self, prs: Presentation) -> None:
        """Slide 3: KPIs Integrales - Dashboard comparativo"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "KPIS INTEGRALES TEMPRANA"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        tf.text = f"COMPARATIVA {self.mes_anterior.upper()} vs {self.mes_actual.upper()}"
        
        # KPIs por cartera
        kpis = self.data.get('kpis_integrales', {})
        
        for cartera, metricas in kpis.items():
            p = tf.add_paragraph()
            p.text = f"\n{cartera.upper()}:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            for kpi, valor in metricas.items():
                p = tf.add_paragraph()
                p.text = f"  • {kpi}: {valor}%"
                p.font.size = Pt(12)
    
    def _create_slide_04_kpis_evolucion(self, prs: Presentation) -> None:
        """Slide 4: Evolución temporal de KPIs"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = f"KPIS CARTERA TEMPRANA - {self.mes_actual.upper()}"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        tf.text = "EVOLUCIÓN DIARIA DE INDICADORES:"
        
        evolucion = self.data.get('kpis_evolucion', {})
        
        if evolucion:
            # Mostrar evolución por días
            for dia, kpis in evolucion.items():
                p = tf.add_paragraph()
                p.text = f"\nDía {dia}:"
                p.font.size = Pt(14)
                p.font.bold = True
                
                for kpi, valor in kpis.items():
                    p = tf.add_paragraph()
                    p.text = f"  • {kpi}: {valor}%"
                    p.font.size = Pt(11)
    
    def _create_slide_05_kpis_contactabilidad(self, prs: Presentation) -> None:
        """Slide 5: Análisis de contactabilidad por servicio"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "KPIS CARTERA TEMPRANA - CONTACTABILIDAD"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        tf.text = "CONTACTABILIDAD POR SERVICIO Y VENCIMIENTO:"
        
        contactabilidad = self.data.get('kpis_contactabilidad', {})
        
        for servicio, datos in contactabilidad.items():
            p = tf.add_paragraph()
            p.text = f"\nSERVICIO {servicio}:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            for vcto, valor in datos.items():
                p = tf.add_paragraph()
                p.text = f"  • {vcto}: {valor}%"
                p.font.size = Pt(12)
    
    def _create_slide_06_kpis_tipos_contacto(self, prs: Presentation) -> None:
        """Slide 6: Tipos de contacto (CD vs CI)"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "KPIS CARTERA TEMPRANA - TIPOS DE CONTACTO"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        tf.text = "CONTACTO DIRECTO vs CONTACTO INDIRECTO:"
        
        tipos_contacto = self.data.get('kpis_tipos_contacto', {})
        
        for servicio, datos in tipos_contacto.items():
            p = tf.add_paragraph()
            p.text = f"\nSERVICIO {servicio}:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            for tipo, valor in datos.items():
                p = tf.add_paragraph()
                p.text = f"  • {tipo}: {valor}%"
                p.font.size = Pt(12)
    
    def _create_slide_07_kpis_resultados(self, prs: Presentation) -> None:
        """Slide 7: KPIs de resultados (Tasa cierre y conversión)"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "KPIS CARTERA TEMPRANA - RESULTADOS"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        tf.text = "TASA DE CIERRE Y CONVERTIBILIDAD:"
        
        resultados = self.data.get('kpis_resultados', {})
        
        for servicio, datos in resultados.items():
            p = tf.add_paragraph()
            p.text = f"\nSERVICIO {servicio}:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            for kpi, valor in datos.items():
                p = tf.add_paragraph()
                p.text = f"  • {kpi}: {valor}%"
                p.font.size = Pt(12)
    
    def _create_slide_08_kpis_esfuerzo(self, prs: Presentation) -> None:
        """Slide 8: KPIs de esfuerzo y efectividad"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "KPIS CARTERA TEMPRANA - ESFUERZO Y EFECTIVIDAD"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        tf.text = "INTENSIDAD Y EFECTIVIDAD:"
        
        esfuerzo = self.data.get('kpis_esfuerzo', {})
        
        for servicio, datos in esfuerzo.items():
            p = tf.add_paragraph()
            p.text = f"\nSERVICIO {servicio}:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            for kpi, valor in datos.items():
                p = tf.add_paragraph()
                p.text = f"  • {kpi}: {valor}"
                p.font.size = Pt(12)
    
    def _create_slide_09_cumplimiento_objetivo(self, prs: Presentation) -> None:
        """Slide 9: Cumplimiento de objetivos"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "CUMPLIMIENTO DE OBJETIVO - TEMPRANA"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        cumplimiento = self.data.get('cumplimiento_objetivo', {})
        cumplimiento_general = cumplimiento.get('cumplimiento_general', 0)
        
        tf.text = f"CUMPLIMIENTO GENERAL: {cumplimiento_general}%"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        
        # Detalle por servicio y vencimiento
        por_servicio = cumplimiento.get('por_servicio', {})
        
        for servicio, vencimientos in por_servicio.items():
            p = tf.add_paragraph()
            p.text = f"\n{servicio}:"
            p.font.size = Pt(16)
            p.font.bold = True
            
            for vcto, datos in vencimientos.items():
                esperado = datos['recupero_esperado']
                real = datos['recupero_real']
                cumplimiento_pct = datos['cumplimiento_%']
                
                p = tf.add_paragraph()
                p.text = f"  • {vcto}: {real}% vs {esperado}% objetivo ({cumplimiento_pct}% cumplimiento)"
                p.font.size = Pt(11)
    
    def _create_slide_10_ranking_agentes(self, prs: Presentation) -> None:
        """Slide 10: Ranking de agentes"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "RANKING DE AGENTES - TEMPRANA"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        tf.text = "TOP 10 AGENTES POR CONVERTIBILIDAD:"
        
        ranking = self.data.get('ranking_agentes', [])
        
        for agente_data in ranking[:10]:
            ranking_pos = agente_data['ranking']
            agente = agente_data['agente']
            convertibilidad = agente_data['convertibilidad_%']
            cuartil = agente_data['cuartil']
            
            p = tf.add_paragraph()
            p.text = f"{ranking_pos}. {agente}: {convertibilidad}% ({cuartil})"
            p.font.size = Pt(12)
    
    def _create_slide_11_estrategia_gestion(self, prs: Presentation) -> None:
        """Slide 11: Estrategia de gestión"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = "ESTRATEGIA DE GESTIÓN"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        tf.text = "METODOLOGÍA DE TRABAJO:"
        
        estrategias = [
            "• Enriquecimiento de datos internos",
            "• Validación y scoring de teléfonos",
            "• Gestión multifono personalizada",
            "• Implementación de intelligence BI",
            "• Automatización con bots de voz",
            "• Seguimiento y optimización continua"
        ]
        
        for estrategia in estrategias:
            p = tf.add_paragraph()
            p.text = estrategia
            p.font.size = Pt(14)
    
    def _create_slide_12_implementacion_bot(self, prs: Presentation) -> None:
        """Slide 12: Implementación de bot"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        title.text = f"IMPLEMENTACIÓN BOT - {self.mes_actual.upper()}"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        tf.text = "AUTOMATIZACIÓN CON BOT DE VOZ:"
        
        bot_features = [
            "• Identificación automática del titular",
            "• Validación de datos personales",
            "• Información de deuda pendiente",
            "• Gestión de promesas de pago",
            "• Escalamiento a agentes humanos",
            "• Registro automático de resultados"
        ]
        
        for feature in bot_features:
            p = tf.add_paragraph()
            p.text = feature
            p.font.size = Pt(14)
        
        # Estadísticas del bot si están disponibles
        p = tf.add_paragraph()
        p.text = f"\nResultados del período:"
        p.font.size = Pt(16)
        p.font.bold = True
        
        # Datos simplificados del voicebot
        p = tf.add_paragraph()
        p.text = f"• Gestiones automatizadas: Variable según período"
        p.font.size = Pt(12)
        
        p = tf.add_paragraph()
        p.text = f"• Tasa de contactabilidad bot: En proceso de optimización"
        p.font.size = Pt(12)
    
    def _create_slide_13_cierre(self, prs: Presentation) -> None:
        """Slide 13: Cierre"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Fondo azul corporativo
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(self.COLORS['movistar_blue'])
        
        # Texto de agradecimiento
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(2)
        thanks_box = slide.shapes.add_textbox(left, top, width, height)
        tf = thanks_box.text_frame
        tf.text = "Gracias!!"
        
        p = tf.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(self.COLORS['white'])
        p.alignment = PP_ALIGN.CENTER
        
        # Logo simulado
        left = Inches(4)
        top = Inches(5.5)
        width = Inches(2)
        height = Inches(1)
        logo_box = slide.shapes.add_textbox(left, top, width, height)
        tf_logo = logo_box.text_frame
        tf_logo.text = "M"
        
        p_logo = tf_logo.paragraphs[0]
        p_logo.font.size = Pt(72)
        p_logo.font.bold = True
        p_logo.font.color.rgb = RGBColor.from_string(self.COLORS['movistar_green'])
        p_logo.alignment = PP_ALIGN.CENTER
    
    def generate_excel_corporate(self, output_path: str = None) -> str:
        """
        Generar Excel corporativo con estructura detallada
        """
        if output_path is None:
            timestamp = self.fecha_generacion.strftime('%Y%m%d_%H%M%S')
            output_path = f"Informe_Detallado_Telefonica_{timestamp}.xlsx"
        
        logger.info(f"Generando Excel corporativo: {output_path}")
        
        # Crear workbook
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        
        # Crear hojas especializadas
        self._create_excel_resumen_ejecutivo_corp(wb)
        self._create_excel_kpis_por_cartera(wb)
        self._create_excel_analisis_servicios(wb)
        self._create_excel_ranking_detallado(wb)
        self._create_excel_cumplimiento_objetivos(wb)
        self._create_excel_data_raw(wb)
        
        # Guardar archivo
        wb.save(output_path)
        logger.info(f"Excel corporativo generado: {output_path}")
        
        return output_path
    
    def _create_excel_resumen_ejecutivo_corp(self, wb: openpyxl.Workbook) -> None:
        """Excel: Hoja de resumen ejecutivo corporativo"""
        ws = wb.create_sheet("Resumen Ejecutivo")
        
        # Título principal
        ws['A1'] = f"INFORME EJECUTIVO GESTIÓN COBRANZA - {self.mes_actual.upper()}"
        ws['A1'].font = Font(bold=True, size=16, color=self.COLORS['white'])
        ws['A1'].fill = PatternFill(start_color=self.COLORS['movistar_blue'], 
                                   end_color=self.COLORS['movistar_blue'], fill_type="solid")
        ws.merge_cells('A1:F1')
        
        # Comparativa mensual
        headers = ['INDICADOR', f'{self.mes_anterior.upper()}', f'{self.mes_actual.upper()}', 'VARIACIÓN', 'ANÁLISIS']
        for i, header in enumerate(headers, 1):
            cell = ws.cell(row=3, column=i, value=header)
            cell.font = Font(bold=True, color=self.COLORS['white'])
            cell.fill = PatternFill(start_color=self.COLORS['movistar_dark_blue'], 
                                   end_color=self.COLORS['movistar_dark_blue'], fill_type="solid")
        
        # Datos de KPIs integrales
        kpis = self.data.get('kpis_integrales', {})
        row = 4
        
        for cartera, metricas in kpis.items():
            for kpi, valor in metricas.items():
                ws.cell(row=row, column=1, value=f"{cartera} - {kpi}")
                ws.cell(row=row, column=2, value="N/A")  # Mes anterior (no disponible)
                ws.cell(row=row, column=3, value=f"{valor}%")
                ws.cell(row=row, column=4, value="N/A")  # Variación
                ws.cell(row=row, column=5, value="En análisis")
                row += 1
        
        # Ajustar anchos
        for col in ['A', 'B', 'C', 'D', 'E']:
            ws.column_dimensions[col].width = 20
    
    def _create_excel_kpis_por_cartera(self, wb: openpyxl.Workbook) -> None:
        """Excel: KPIs detallados por cartera"""
        ws = wb.create_sheet("KPIs por Cartera")
        
        ws['A1'] = "KPIS DETALLADOS POR CARTERA"
        ws['A1'].font = Font(bold=True, size=14, color=self.COLORS['white'])
        ws['A1'].fill = PatternFill(start_color=self.COLORS['movistar_blue'], 
                                   end_color=self.COLORS['movistar_blue'], fill_type="solid")
        ws.merge_cells('A1:H1')
        
        # Headers
        headers = ['CARTERA', 'CONTACTABILIDAD%', 'CONTACTO_DIRECTO%', 'CONTACTO_INDIRECTO%', 
                  'CONVERSIÓN%', 'INTENSIDAD', 'EFECTIVIDAD%', 'OBSERVACIONES']
        
        for i, header in enumerate(headers, 1):
            cell = ws.cell(row=3, column=i, value=header)
            cell.font = Font(bold=True, color=self.COLORS['white'])
            cell.fill = PatternFill(start_color=self.COLORS['telefonica_orange'], 
                                   end_color=self.COLORS['telefonica_orange'], fill_type="solid")
        
        # Datos
        kpis = self.data.get('kpis_integrales', {})
        row = 4
        
        for cartera, metricas in kpis.items():
            ws.cell(row=row, column=1, value=cartera)
            ws.cell(row=row, column=2, value=metricas.get('%CONT', 0))
            ws.cell(row=row, column=3, value=metricas.get('CD%', 0))
            ws.cell(row=row, column=4, value=metricas.get('CI%', 0))
            ws.cell(row=row, column=5, value=metricas.get('%CONV', 0))
            ws.cell(row=row, column=6, value=metricas.get('INTENSIDAD', 0))
            ws.cell(row=row, column=7, value=metricas.get('EFECTIVIDAD%', 0))
            ws.cell(row=row, column=8, value="Análisis en progreso")
            row += 1
        
        # Ajustar anchos
        for col in ['A', 'B', 'C', 'D', 'E', 'F', 'G', 'H']:
            ws.column_dimensions[col].width = 15
    
    def _create_excel_analisis_servicios(self, wb: openpyxl.Workbook) -> None:
        """Excel: Análisis por servicios FIJA vs MÓVIL"""
        ws = wb.create_sheet("Análisis Servicios")
        
        ws['A1'] = "ANÁLISIS DETALLADO POR SERVICIO"
        ws['A1'].font = Font(bold=True, size=14, color=self.COLORS['white'])
        ws['A1'].fill = PatternFill(start_color=self.COLORS['movistar_green'], 
                                   end_color=self.COLORS['movistar_green'], fill_type="solid")
        ws.merge_cells('A1:F1')
        
        # Datos de contactabilidad por servicio
        contactabilidad = self.data.get('kpis_contactabilidad', {})
        
        row = 3
        for servicio, datos in contactabilidad.items():
            ws.cell(row=row, column=1, value=f"SERVICIO {servicio}")
            ws.cell(row=row, column=1).font = Font(bold=True)
            row += 1
            
            for vcto, valor in datos.items():
                ws.cell(row=row, column=2, value=vcto)
                ws.cell(row=row, column=3, value=f"{valor}%")
                row += 1
            
            row += 1  # Espacio entre servicios
    
    def _create_excel_ranking_detallado(self, wb: openpyxl.Workbook) -> None:
        """Excel: Ranking detallado de agentes"""
        ws = wb.create_sheet("Ranking Agentes")
        
        ws['A1'] = "RANKING DETALLADO DE AGENTES"
        ws['A1'].font = Font(bold=True, size=14, color=self.COLORS['white'])
        ws['A1'].fill = PatternFill(start_color=self.COLORS['telefonica_blue'], 
                                   end_color=self.COLORS['telefonica_blue'], fill_type="solid")
        ws.merge_cells('A1:H1')
        
        # Headers
        headers = ['RANKING', 'AGENTE', 'GESTIONES', 'CEF', 'NEF', 'MONTO_PAGADO', 
                  'CONVERTIBILIDAD%', 'CUARTIL']
        
        for i, header in enumerate(headers, 1):
            cell = ws.cell(row=3, column=i, value=header)
            cell.font = Font(bold=True, color=self.COLORS['white'])
            cell.fill = PatternFill(start_color=self.COLORS['telefonica_light_blue'], 
                                   end_color=self.COLORS['telefonica_light_blue'], fill_type="solid")
        
        # Datos del ranking
        ranking = self.data.get('ranking_agentes', [])
        
        for i, agente_data in enumerate(ranking, 4):
            ws.cell(row=i, column=1, value=agente_data['ranking'])
            ws.cell(row=i, column=2, value=agente_data['agente'])
            ws.cell(row=i, column=3, value=agente_data['gestiones'])
            ws.cell(row=i, column=4, value=agente_data['cef'])
            ws.cell(row=i, column=5, value=agente_data['nef'])
            ws.cell(row=i, column=6, value=agente_data['monto_pagado'])
            ws.cell(row=i, column=7, value=agente_data['convertibilidad_%'])
            ws.cell(row=i, column=8, value=agente_data['cuartil'])
        
        # Ajustar anchos
        for col in ['A', 'B', 'C', 'D', 'E', 'F', 'G', 'H']:
            ws.column_dimensions[col].width = 15
    
    def _create_excel_cumplimiento_objetivos(self, wb: openpyxl.Workbook) -> None:
        """Excel: Cumplimiento de objetivos detallado"""
        ws = wb.create_sheet("Cumplimiento Objetivos")
        
        ws['A1'] = "CUMPLIMIENTO DETALLADO DE OBJETIVOS"
        ws['A1'].font = Font(bold=True, size=14, color=self.COLORS['white'])
        ws['A1'].fill = PatternFill(start_color=self.COLORS['movistar_blue'], 
                                   end_color=self.COLORS['movistar_blue'], fill_type="solid")
        ws.merge_cells('A1:F1')
        
        # Cumplimiento general
        cumplimiento = self.data.get('cumplimiento_objetivo', {})
        cumplimiento_general = cumplimiento.get('cumplimiento_general', 0)
        
        ws['A3'] = "CUMPLIMIENTO GENERAL:"
        ws['B3'] = f"{cumplimiento_general}%"
        ws['A3'].font = Font(bold=True, size=12)
        ws['B3'].font = Font(bold=True, size=12, color='FF0000' if cumplimiento_general < 70 else '00AA00')
        
        # Headers detalle
        headers = ['SERVICIO', 'VENCIMIENTO', 'OBJETIVO%', 'REAL%', 'CUMPLIMIENTO%', 'PESO']
        
        for i, header in enumerate(headers, 1):
            cell = ws.cell(row=5, column=i, value=header)
            cell.font = Font(bold=True, color=self.COLORS['white'])
            cell.fill = PatternFill(start_color=self.COLORS['telefonica_orange'], 
                                   end_color=self.COLORS['telefonica_orange'], fill_type="solid")
        
        # Datos detallados
        por_servicio = cumplimiento.get('por_servicio', {})
        row = 6
        
        for servicio, vencimientos in por_servicio.items():
            for vcto, datos in vencimientos.items():
                ws.cell(row=row, column=1, value=servicio)
                ws.cell(row=row, column=2, value=vcto)
                ws.cell(row=row, column=3, value=datos['recupero_esperado'])
                ws.cell(row=row, column=4, value=datos['recupero_real'])
                ws.cell(row=row, column=5, value=datos['cumplimiento_%'])
                ws.cell(row=row, column=6, value=round(datos['peso'], 0))
                row += 1
    
    def _create_excel_data_raw(self, wb: openpyxl.Workbook) -> None:
        """Excel: Datos raw para análisis"""
        ws = wb.create_sheet("Datos Raw")
        
        ws['A1'] = "DATOS CONSOLIDADOS PARA ANÁLISIS"
        ws['A1'].font = Font(bold=True, size=14)
        
        # Estructura de datos disponibles
        estructura = [
            "ESTRUCTURA DE DATOS DISPONIBLES:",
            "",
            "1. Asignación por Cartera:",
            f"   - {len(self.data.get('asignacion_cartera', {}).get('resumen_cartera', {}))} carteras analizadas",
            "",
            "2. KPIs Integrales:",
            f"   - {len(self.data.get('kpis_integrales', {}))} carteras con KPIs",
            "",
            "3. Ranking de Agentes:",
            f"   - {len(self.data.get('ranking_agentes', []))} agentes evaluados",
            "",
            "4. Cumplimiento de Objetivos:",
            f"   - Cumplimiento general: {self.data.get('cumplimiento_objetivo', {}).get('cumplimiento_general', 0)}%",
            "",
            "5. Período Analizado:",
            f"   - Inicio: {self.fecha_inicio}",
            f"   - Fin: {self.fecha_fin}",
            f"   - Mes actual: {self.mes_actual}",
            f"   - Mes anterior: {self.mes_anterior}"
        ]
        
        for i, item in enumerate(estructura, 3):
            ws.cell(row=i, column=1, value=item)
            if item.startswith(("1.", "2.", "3.", "4.", "5.")):
                ws.cell(row=i, column=1).font = Font(bold=True)
    
    def generate_complete_corporate_report(self, output_dir: str = None) -> Tuple[str, str]:
        """
        Generar reporte corporativo completo (Excel + PowerPoint)
        """
        if output_dir is None:
            output_dir = tempfile.gettempdir()
        
        timestamp = self.fecha_generacion.strftime('%Y%m%d_%H%M%S')
        
        excel_path = os.path.join(output_dir, f"Informe_Detallado_Telefonica_{timestamp}.xlsx")
        ppt_path = os.path.join(output_dir, f"Informe_Gestion_Cobranza_Telefonica_{timestamp}.pptx")
        
        # Generar ambos reportes corporativos
        excel_file = self.generate_excel_corporate(excel_path)
        ppt_file = self.generate_powerpoint_corporate(ppt_path)
        
        logger.info(f"Reportes corporativos generados en: {output_dir}")
        
        return excel_file, ppt_file
