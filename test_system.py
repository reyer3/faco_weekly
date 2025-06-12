#!/usr/bin/env python3
"""
Test Script - FACO Weekly Sistema Corregido
==========================================

Script de prueba para verificar que el sistema funciona correctamente
sin errores de compilación y con todas las dependencias instaladas.
"""

import sys
import os
import traceback
from datetime import datetime, timedelta

def test_basic_imports():
    """Test 1: Verificar imports básicos"""
    print("🔍 Test 1: Verificando imports básicos...")
    
    try:
        import fastapi
        print("  ✅ FastAPI: OK")
    except ImportError as e:
        print(f"  ❌ FastAPI: ERROR - {e}")
        return False
    
    try:
        import pandas as pd
        print("  ✅ Pandas: OK")
    except ImportError as e:
        print(f"  ❌ Pandas: ERROR - {e}")
        return False
    
    try:
        import openpyxl
        print("  ✅ OpenPyXL: OK")
    except ImportError as e:
        print(f"  ❌ OpenPyXL: ERROR - {e}")
        return False
    
    try:
        import pptx
        print("  ✅ Python-PPTX: OK")
    except ImportError as e:
        print(f"  ❌ Python-PPTX: ERROR - {e}")
        return False
    
    try:
        from google.cloud import bigquery
        print("  ✅ BigQuery: OK")
    except ImportError as e:
        print(f"  ⚠️ BigQuery: ERROR - {e} (requiere credenciales)")
    
    return True

def test_report_generator():
    """Test 2: Verificar generador de reportes"""
    print("\n🔍 Test 2: Verificando generador de reportes...")
    
    try:
        from report_generator import TelefonicaReportGenerator
        print("  ✅ Import TelefonicaReportGenerator: OK")
        
        # Crear instancia
        generator = TelefonicaReportGenerator("2025-06-01", "2025-06-12")
        print("  ✅ Instancia creada: OK")
        
        # Verificar estructura de datos
        assert hasattr(generator, 'data'), "Falta atributo 'data'"
        assert hasattr(generator, 'COLORS'), "Falta atributo 'COLORS'"
        print("  ✅ Estructura de datos: OK")
        
        return True
        
    except Exception as e:
        print(f"  ❌ Error en generador: {e}")
        print(f"  📝 Traceback: {traceback.format_exc()}")
        return False

def test_main_api():
    """Test 3: Verificar API principal"""
    print("\n🔍 Test 3: Verificando API principal...")
    
    try:
        import main
        print("  ✅ Import main: OK")
        
        # Verificar que la app existe
        assert hasattr(main, 'app'), "Falta objeto 'app' en main.py"
        print("  ✅ FastAPI app: OK")
        
        return True
        
    except Exception as e:
        print(f"  ❌ Error en main: {e}")
        print(f"  📝 Traceback: {traceback.format_exc()}")
        return False

def test_dummy_excel_generation():
    """Test 4: Verificar generación de Excel dummy"""
    print("\n🔍 Test 4: Verificando generación de Excel dummy...")
    
    try:
        import openpyxl
        import tempfile
        import os
        
        # Crear archivo Excel de prueba
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Test Sheet"
        
        # Agregar algunos datos
        ws['A1'] = "FACO Weekly Test"
        ws['A2'] = f"Generado: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        ws['A3'] = "Estado: Funcionando correctamente"
        
        # Guardar en directorio temporal
        temp_dir = tempfile.gettempdir()
        test_file = os.path.join(temp_dir, "faco_test.xlsx")
        wb.save(test_file)
        
        # Verificar que existe
        if os.path.exists(test_file):
            file_size = os.path.getsize(test_file)
            print(f"  ✅ Excel generado: {test_file} ({file_size} bytes)")
            
            # Limpiar archivo de prueba
            os.remove(test_file)
            return True
        else:
            print("  ❌ No se pudo generar el archivo Excel")
            return False
            
    except Exception as e:
        print(f"  ❌ Error generando Excel: {e}")
        return False

def test_dummy_powerpoint_generation():
    """Test 5: Verificar generación de PowerPoint dummy"""
    print("\n🔍 Test 5: Verificando generación de PowerPoint dummy...")
    
    try:
        from pptx import Presentation
        import tempfile
        import os
        
        # Crear presentación de prueba
        prs = Presentation()
        
        # Slide de prueba
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "FACO Weekly Test"
        subtitle.text = f"Sistema funcionando correctamente\\n{datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        # Guardar en directorio temporal
        temp_dir = tempfile.gettempdir()
        test_file = os.path.join(temp_dir, "faco_test.pptx")
        prs.save(test_file)
        
        # Verificar que existe
        if os.path.exists(test_file):
            file_size = os.path.getsize(test_file)
            print(f"  ✅ PowerPoint generado: {test_file} ({file_size} bytes)")
            
            # Limpiar archivo de prueba
            os.remove(test_file)
            return True
        else:
            print("  ❌ No se pudo generar el archivo PowerPoint")
            return False
            
    except Exception as e:
        print(f"  ❌ Error generando PowerPoint: {e}")
        return False

def test_data_processing():
    """Test 6: Verificar procesamiento de datos dummy"""
    print("\n🔍 Test 6: Verificando procesamiento de datos dummy...")
    
    try:
        import pandas as pd
        import numpy as np
        
        # Crear datos dummy para prueba
        gestiones_dummy = pd.DataFrame({
            'canal': ['CALL', 'VOICEBOT', 'CALL', 'VOICEBOT'] * 100,
            'contactabilidad': ['CONTACTO_EFECTIVO', 'NO_CONTACTO', 'CONTACTO_EFECTIVO', 'CONTACTO_NO_EFECTIVO'] * 100,
            'es_pdp': ['SI', 'NO', 'SI', 'NO'] * 100,
            'cod_luna': range(400),
            'monto_compromiso': np.random.uniform(10, 1000, 400),
            'date': pd.date_range('2025-06-01', periods=400, freq='H'),
            'duracion': np.random.uniform(30, 180, 400)
        })
        
        print(f"  ✅ DataFrame dummy creado: {len(gestiones_dummy)} registros")
        
        # Procesar datos básicos
        call_data = gestiones_dummy[gestiones_dummy['canal'] == 'CALL']
        voicebot_data = gestiones_dummy[gestiones_dummy['canal'] == 'VOICEBOT']
        
        call_contactos = len(call_data[call_data['contactabilidad'] == 'CONTACTO_EFECTIVO'])
        voicebot_contactos = len(voicebot_data[voicebot_data['contactabilidad'] == 'CONTACTO_EFECTIVO'])
        
        print(f"  ✅ CALL contactos efectivos: {call_contactos}")
        print(f"  ✅ VOICEBOT contactos efectivos: {voicebot_contactos}")
        
        # Calcular tasas
        tasa_call = round(call_contactos / len(call_data) * 100, 2) if len(call_data) > 0 else 0
        tasa_voicebot = round(voicebot_contactos / len(voicebot_data) * 100, 2) if len(voicebot_data) > 0 else 0
        
        print(f"  ✅ Tasa contactabilidad CALL: {tasa_call}%")
        print(f"  ✅ Tasa contactabilidad VOICEBOT: {tasa_voicebot}%")
        
        return True
        
    except Exception as e:
        print(f"  ❌ Error procesando datos: {e}")
        return False

def main():
    """Ejecutar todos los tests"""
    print("=" * 60)
    print("🧪 FACO WEEKLY - SUITE DE PRUEBAS")
    print("Verificando sistema sin errores de compilación")
    print("=" * 60)
    
    tests = [
        test_basic_imports,
        test_report_generator,
        test_main_api,
        test_dummy_excel_generation,
        test_dummy_powerpoint_generation,
        test_data_processing
    ]
    
    results = []
    
    for test in tests:
        try:
            result = test()
            results.append(result)
        except Exception as e:
            print(f"❌ Test falló con excepción: {e}")
            results.append(False)
    
    # Resumen final
    print("\n" + "=" * 60)
    print("📊 RESUMEN DE PRUEBAS")
    print("=" * 60)
    
    passed = sum(results)
    total = len(results)
    
    for i, (test, result) in enumerate(zip(tests, results), 1):
        status = "✅ PASSED" if result else "❌ FAILED"
        print(f"Test {i}: {test.__name__} - {status}")
    
    print(f"\n🎯 Resultado: {passed}/{total} pruebas exitosas")
    
    if passed == total:
        print("🎉 ¡TODOS LOS TESTS PASARON!")
        print("✅ Sistema listo para generar reportes automatizados")
        print("\n📋 Próximos pasos:")
        print("1. Configurar credenciales Google Cloud")
        print("2. Ejecutar: python3 main.py")
        print("3. Probar: curl http://localhost:8000/health")
        print("4. Generar reporte: curl -X POST http://localhost:8000/generate-reports \\")
        print('   -H "Content-Type: application/json" \\')
        print('   -d \'{"fecha_inicio": "2025-06-01", "fecha_fin": "2025-06-12"}\'')
        return True
    else:
        print("⚠️ Algunos tests fallaron")
        print("🔧 Revisa la instalación siguiendo README.md")
        print("📞 Contacta soporte si persisten los errores")
        return False

if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)
