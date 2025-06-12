"""
REPORT GENERATOR - TelefÃ³nica del PerÃº (VERSIÃ“N CORREGIDA)
===========================================================

MÃ³dulo para generar reportes semanales automatizados en Excel y PowerPoint
SIN dependencias problemÃ¡ticas de matplotlib o compilaciÃ³n C.

CaracterÃ­sticas:
- GeneraciÃ³n Excel con openpyxl (sin matplotlib)
- Presentaciones PowerPoint con python-pptx
- AnÃ¡lisis consolidado CALL vs VOICEBOT
- MÃ©tricas ejecutivas automatizadas
- Formato corporativo TelefÃ³nica

Autor: Sistema Automatizado FACO Weekly
Fecha: Junio 2025 - VersiÃ³n Corregida
"""

import pandas as pd
import numpy as np
from datetime import datetime, date, timedelta
import tempfile
import os
from typing import Dict, List, Optional, Tuple
import logging

# Excel libraries (sin matplotlib)
import openpyxl
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.worksheet.table import Table, TableStyleInfo
import xlsxwriter

# PowerPoint libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

class TelefonicaReportGenerator:
    """
    Generador de reportes semanales para TelefÃ³nica del PerÃº
    VERSIÃ“N CORREGIDA - Sin dependencias problemÃ¡ticas
    """
    
    # Colores corporativos TelefÃ³nica
    COLORS = {
        'telefonica_blue': '0019A5',
        'telefonica_light_blue': '5BB4E5', 
        'telefonica_green': '00A651',
        'telefonica_orange': 'FF6600',
        'telefonica_gray': '666666',
        'white': 'FFFFFF',
        'light_gray': 'F2F2F2'
    }
    
    def __init__(self, fecha_inicio: str, fecha_fin: str):
        """
        Inicializar generador de reportes
        
        Args:
            fecha_inicio: Fecha inicio perÃ­odo (YYYY-MM-DD)
            fecha_fin: Fecha fin perÃ­odo (YYYY-MM-DD)
        """
        self.fecha_inicio = fecha_inicio
        self.fecha_fin = fecha_fin
        self.periodo_str = f"{fecha_inicio} a {fecha_fin}"
        self.fecha_generacion = datetime.now()
        
        # Estructura de datos
        self.data = {
            'resumen_ejecutivo': {},
            'canal_call': {},
            'canal_voicebot': {},
            'evolucion_diaria': [],
            'carteras_activas': [],
            'kpis_consolidados': {},
            'recomendaciones': []
        }
        
        logger.info(f"ðŸš€ Iniciando generador de reportes para perÃ­odo: {self.periodo_str}")
    
    def load_data_from_processing(self, 
                                gestiones_df: pd.DataFrame,
                                calendario_df: pd.DataFrame,
                                asignacion_df: pd.DataFrame,
                                pagos_df: pd.DataFrame,
                                kpis_campania: List[Dict]) -> None:
        """
        Cargar datos desde el procesamiento principal
        """
        logger.info("ðŸ“Š Cargando datos para generaciÃ³n de reportes")
        
        try:
            # Procesar datos de gestiones
            self._process_gestiones_data(gestiones_df)
            
            # Procesar datos de calendario y asignaciones
            self._process_calendario_data(calendario_df, asignacion_df)
            
            # Procesar datos de pagos
            self._process_pagos_data(pagos_df)
            
            # Procesar KPIs por campaÃ±a
            self._process_kpis_campania(kpis_campania)
            
            # Calcular mÃ©tricas consolidadas
            self._calculate_consolidated_metrics()
            
            # Generar recomendaciones automÃ¡ticas
            self._generate_recommendations()
            
            logger.info("âœ… Datos cargados y procesados exitosamente")
            
        except Exception as e:
            logger.error(f"âŒ Error cargando datos: {str(e)}")
            raise
    
    def _process_gestiones_data(self, gestiones_df: pd.DataFrame) -> None:
        """Procesar datos de gestiones por canal"""
        if gestiones_df.empty:
            logger.warning("âš ï¸ No hay datos de gestiones para procesar")
            return
        
        try:
            # Separar por canal
            call_data = gestiones_df[gestiones_df['canal'] == 'CALL']
            voicebot_data = gestiones_df[gestiones_df['canal'] == 'VOICEBOT']
            
            # MÃ©tricas CALL
            self.data['canal_call'] = self._calculate_channel_metrics(call_data, 'CALL')
            
            # MÃ©tricas VOICEBOT
            self.data['canal_voicebot'] = self._calculate_channel_metrics(voicebot_data, 'VOICEBOT')
            
            # EvoluciÃ³n diaria
            self._calculate_daily_evolution(gestiones_df)
            
        except Exception as e:
            logger.error(f"Error procesando gestiones: {str(e)}")
            raise
    
    def _calculate_channel_metrics(self, channel_data: pd.DataFrame, channel_name: str) -> Dict:
        """Calcular mÃ©tricas para un canal especÃ­fico"""
        if channel_data.empty:
            return {
                'total_gestiones': 0,
                'contactos_efectivos': 0,
                'contactos_no_efectivos': 0,
                'no_contactos': 0,
                'compromisos': 0,
                'monto_compromisos': 0,
                'clientes_unicos': 0,
                'tasa_contactabilidad': 0,
                'tasa_compromiso': 0,
                'duracion_promedio': 0
            }
        
        # MÃ©tricas bÃ¡sicas
        total_gestiones = len(channel_data)
        contactos_efectivos = len(channel_data[channel_data['contactabilidad'] == 'CONTACTO_EFECTIVO'])
        contactos_no_efectivos = len(channel_data[channel_data['contactabilidad'] == 'CONTACTO_NO_EFECTIVO'])
        no_contactos = len(channel_data[channel_data['contactabilidad'] == 'NO_CONTACTO'])
        compromisos = len(channel_data[channel_data['es_pdp'] == 'SI'])
        monto_compromisos = channel_data['monto_compromiso'].sum() if 'monto_compromiso' in channel_data.columns else 0
        clientes_unicos = channel_data['cod_luna'].nunique()
        
        # DuraciÃ³n promedio (solo para CALL)
        duracion_promedio = 0
        if channel_name == 'CALL' and 'duracion' in channel_data.columns:
            duracion_promedio = channel_data['duracion'].mean()
        
        # Calcular tasas
        tasa_contactabilidad = round(contactos_efectivos / max(total_gestiones, 1) * 100, 2)
        tasa_compromiso = round(compromisos / max(contactos_efectivos, 1) * 100, 2)
        
        return {
            'total_gestiones': total_gestiones,
            'contactos_efectivos': contactos_efectivos,
            'contactos_no_efectivos': contactos_no_efectivos,
            'no_contactos': no_contactos,
            'compromisos': compromisos,
            'monto_compromisos': monto_compromisos,
            'clientes_unicos': clientes_unicos,
            'tasa_contactabilidad': tasa_contactabilidad,
            'tasa_compromiso': tasa_compromiso,
            'duracion_promedio': round(duracion_promedio, 1)
        }
    
    def _calculate_daily_evolution(self, gestiones_df: pd.DataFrame) -> None:
        """Calcular evoluciÃ³n diaria"""
        try:
            if gestiones_df.empty:
                return
            
            # Asegurar que date es datetime
            if 'date' in gestiones_df.columns:
                gestiones_df['fecha'] = pd.to_datetime(gestiones_df['date']).dt.date
            else:
                logger.warning("No hay columna 'date' en gestiones")
                return
            
            # Separar por canal
            call_data = gestiones_df[gestiones_df['canal'] == 'CALL']
            voicebot_data = gestiones_df[gestiones_df['canal'] == 'VOICEBOT']
            
            # EvoluciÃ³n CALL
            evolucion_call = call_data.groupby('fecha').agg({
                'cod_luna': 'count',
                'contactabilidad': lambda x: (x == 'CONTACTO_EFECTIVO').sum()
            }).rename(columns={'cod_luna': 'gestiones_call', 'contactabilidad': 'contactos_call'})
            
            # EvoluciÃ³n VOICEBOT
            evolucion_voicebot = voicebot_data.groupby('fecha').agg({
                'cod_luna': 'count',
                'contactabilidad': lambda x: (x == 'CONTACTO_EFECTIVO').sum()
            }).rename(columns={'cod_luna': 'gestiones_voicebot', 'contactabilidad': 'contactos_voicebot'})
            
            # Combinar evoluciÃ³n diaria
            evolucion_completa = evolucion_call.join(evolucion_voicebot, how='outer').fillna(0)
            evolucion_completa['total_gestiones'] = evolucion_completa['gestiones_call'] + evolucion_completa['gestiones_voicebot']
            evolucion_completa['total_contactos'] = evolucion_completa['contactos_call'] + evolucion_completa['contactos_voicebot']
            
            self.data['evolucion_diaria'] = [
                {
                    'fecha': fecha.strftime('%Y-%m-%d'),
                    'call_gestiones': int(row['gestiones_call']),
                    'call_contactos': int(row['contactos_call']),
                    'voicebot_gestiones': int(row['gestiones_voicebot']),
                    'voicebot_contactos': int(row['contactos_voicebot']),
                    'total_gestiones': int(row['total_gestiones']),
                    'total_contactos': int(row['total_contactos']),
                    'tasa_contactabilidad': round(row['total_contactos'] / max(row['total_gestiones'], 1) * 100, 2)
                }
                for fecha, row in evolucion_completa.iterrows()
            ]
            
        except Exception as e:
            logger.error(f"Error calculando evoluciÃ³n diaria: {str(e)}")
            self.data['evolucion_diaria'] = []
    
    def _process_calendario_data(self, calendario_df: pd.DataFrame, asignacion_df: pd.DataFrame) -> None:
        """Procesar datos de calendario y asignaciones"""
        try:
            if not calendario_df.empty:
                self.data['carteras_activas'] = [
                    {
                        'archivo': row['archivo'],
                        'tipo_cartera': row.get('tipo_cartera', 'N/A'),
                        'fecha_asignacion': row['fecha_asignacion'].strftime('%Y-%m-%d') if hasattr(row['fecha_asignacion'], 'strftime') else str(row['fecha_asignacion']),
                        'fecha_cierre': row['fecha_cierre'].strftime('%Y-%m-%d') if hasattr(row['fecha_cierre'], 'strftime') else str(row['fecha_cierre']),
                        'suma_lineas': row.get('suma_lineas', 0),
                        'dias_vigencia': row.get('dias_vigencia', 0),
                        'estado': row.get('estado_vigencia', 'ACTIVA')
                    }
                    for _, row in calendario_df.iterrows()
                ]
            
            # Agregar informaciÃ³n de asignaciones
            if not asignacion_df.empty:
                for cartera in self.data['carteras_activas']:
                    asig_data = asignacion_df[asignacion_df['archivo'] == cartera['archivo']]
                    if not asig_data.empty:
                        cartera.update({
                            'clientes_asignados': int(asig_data.iloc[0].get('clientes_asignados', 0)),
                            'cuentas_asignadas': int(asig_data.iloc[0].get('cuentas_asignadas', 0))
                        })
                        
        except Exception as e:
            logger.error(f"Error procesando calendario: {str(e)}")
            self.data['carteras_activas'] = []
    
    def _process_pagos_data(self, pagos_df: pd.DataFrame) -> None:
        """Procesar datos de pagos"""
        try:
            if not pagos_df.empty:
                self.data['pagos'] = {
                    'total_pagos': len(pagos_df),
                    'clientes_con_pago': pagos_df['nro_documento'].nunique(),
                    'monto_total': pagos_df['monto_cancelado'].sum(),
                    'ticket_promedio': pagos_df['monto_cancelado'].mean(),
                    'monto_min': pagos_df['monto_cancelado'].min(),
                    'monto_max': pagos_df['monto_cancelado'].max()
                }
            else:
                self.data['pagos'] = {
                    'total_pagos': 0,
                    'clientes_con_pago': 0,
                    'monto_total': 0,
                    'ticket_promedio': 0,
                    'monto_min': 0,
                    'monto_max': 0
                }
        except Exception as e:
            logger.error(f"Error procesando pagos: {str(e)}")
            self.data['pagos'] = {'total_pagos': 0, 'clientes_con_pago': 0, 'monto_total': 0, 'ticket_promedio': 0, 'monto_min': 0, 'monto_max': 0}
    
    def _process_kpis_campania(self, kpis_campania: List[Dict]) -> None:
        """Procesar KPIs por campaÃ±a"""
        self.data['kpis_por_campania'] = kpis_campania or []
    
    def _calculate_consolidated_metrics(self) -> None:
        """Calcular mÃ©tricas consolidadas"""
        try:
            call = self.data['canal_call']
            voicebot = self.data['canal_voicebot']
            
            total_gestiones = call.get('total_gestiones', 0) + voicebot.get('total_gestiones', 0)
            total_contactos = call.get('contactos_efectivos', 0) + voicebot.get('contactos_efectivos', 0)
            total_compromisos = call.get('compromisos', 0) + voicebot.get('compromisos', 0)
            
            self.data['resumen_ejecutivo'] = {
                'total_gestiones': total_gestiones,
                'total_contactos_efectivos': total_contactos,
                'total_compromisos': total_compromisos,
                'tasa_contactabilidad_global': round(total_contactos / max(total_gestiones, 1) * 100, 2),
                'tasa_compromiso_global': round(total_compromisos / max(total_contactos, 1) * 100, 2),
                'monto_compromisos_call': call.get('monto_compromisos', 0),
                'clientes_unicos_total': call.get('clientes_unicos', 0) + voicebot.get('clientes_unicos', 0)
            }
            
        except Exception as e:
            logger.error(f"Error calculando mÃ©tricas consolidadas: {str(e)}")
            self.data['resumen_ejecutivo'] = {
                'total_gestiones': 0,
                'total_contactos_efectivos': 0,
                'total_compromisos': 0,
                'tasa_contactabilidad_global': 0,
                'tasa_compromiso_global': 0,
                'monto_compromisos_call': 0,
                'clientes_unicos_total': 0
            }
    
    def _generate_recommendations(self) -> None:
        """Generar recomendaciones automÃ¡ticas basadas en datos"""
        recomendaciones = []
        
        try:
            # Analizar contactabilidad por canal
            call_contactabilidad = self.data['canal_call'].get('tasa_contactabilidad', 0)
            voicebot_contactabilidad = self.data['canal_voicebot'].get('tasa_contactabilidad', 0)
            
            if voicebot_contactabilidad < 2.0:
                recomendaciones.append({
                    'categoria': 'OptimizaciÃ³n VOICEBOT',
                    'prioridad': 'Alta',
                    'descripcion': f'Mejorar scripts VOICEBOT - actual: {voicebot_contactabilidad}%, meta: 2%+',
                    'accion': 'Revisar y optimizar scripts de contacto automatizado'
                })
            
            if call_contactabilidad > voicebot_contactabilidad * 3:
                recomendaciones.append({
                    'categoria': 'Balanceo de Canales',
                    'prioridad': 'Media',
                    'descripcion': f'CALL ({call_contactabilidad}%) vs VOICEBOT ({voicebot_contactabilidad}%) - gran diferencia',
                    'accion': 'Redistribuir cartera para optimizar efectividad'
                })
            
            # Analizar compromisos
            monto_compromisos = self.data['canal_call'].get('monto_compromisos', 0)
            if monto_compromisos > 100000:
                recomendaciones.append({
                    'categoria': 'Seguimiento Compromisos',
                    'prioridad': 'Alta',
                    'descripcion': f'${monto_compromisos:,.0f} en compromisos requiere seguimiento intensivo',
                    'accion': 'Implementar sistema de tracking de cumplimiento'
                })
            
            self.data['recomendaciones'] = recomendaciones
            
        except Exception as e:
            logger.error(f"Error generando recomendaciones: {str(e)}")
            self.data['recomendaciones'] = []
    
    def generate_excel_report(self, output_path: str = None) -> str:
        """
        Generar reporte Excel completo
        """
        if output_path is None:
            timestamp = self.fecha_generacion.strftime('%Y%m%d_%H%M%S')
            output_path = f"Informe_Semanal_Telefonica_{timestamp}.xlsx"
        
        logger.info(f"ðŸ“Š Generando reporte Excel: {output_path}")
        
        try:
            # Crear workbook
            wb = openpyxl.Workbook()
            wb.remove(wb.active)  # Remover hoja por defecto
            
            # Generar hojas
            self._create_excel_resumen_ejecutivo(wb)
            self._create_excel_analisis_canales(wb)
            self._create_excel_evolucion_diaria(wb)
            self._create_excel_carteras_activas(wb)
            if self.data['kpis_por_campania']:
                self._create_excel_kpis_campanias(wb)
            if self.data['recomendaciones']:
                self._create_excel_recomendaciones(wb)
            
            # Guardar archivo
            wb.save(output_path)
            logger.info(f"âœ… Reporte Excel generado exitosamente: {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"âŒ Error generando Excel: {str(e)}")
            raise
    
    def _create_excel_resumen_ejecutivo(self, wb: openpyxl.Workbook) -> None:
        """Crear hoja de resumen ejecutivo"""
        ws = wb.create_sheet("Resumen Ejecutivo")
        
        # TÃ­tulo principal
        ws['A1'] = "INFORME SEMANAL DE GESTIÃ“N DE COBRANZA"
        ws['A2'] = f"TelefÃ³nica del PerÃº - PerÃ­odo: {self.periodo_str}"
        ws['A3'] = f"Generado: {self.fecha_generacion.strftime('%d/%m/%Y %H:%M')}"
        
        # Formatear tÃ­tulos
        for row in range(1, 4):
            cell = ws[f'A{row}']
            cell.font = Font(bold=True, size=14 if row == 1 else 12, color=self.COLORS['white'])
            cell.fill = PatternFill(start_color=self.COLORS['telefonica_blue'], 
                                  end_color=self.COLORS['telefonica_blue'], fill_type="solid")
            ws.merge_cells(f'A{row}:D{row}')
        
        # Datos principales
        resumen = self.data['resumen_ejecutivo']
        data_rows = [
            ['', '', '', ''],
            ['INDICADOR CLAVE', 'VALOR', 'MÃ‰TRICA', 'OBSERVACIONES'],
            ['Total Gestiones', f"{resumen.get('total_gestiones', 0):,}", '100%', 'CALL + VOICEBOT'],
            ['Contactos Efectivos', f"{resumen.get('total_contactos_efectivos', 0):,}", 
             f"{resumen.get('tasa_contactabilidad_global', 0)}%", 'Tasa de contactabilidad global'],
            ['Compromisos Obtenidos', f"{resumen.get('total_compromisos', 0):,}", 
             f"{resumen.get('tasa_compromiso_global', 0)}%", 'De contactos efectivos'],
            ['Monto Compromisos CALL', f"${resumen.get('monto_compromisos_call', 0):,.0f}", '-', 
             f"Promedio: ${resumen.get('monto_compromisos_call', 0) / max(resumen.get('total_compromisos', 1), 1):.0f}"],
            ['Clientes Ãšnicos', f"{resumen.get('clientes_unicos_total', 0):,}", '-', 'Total gestionados'],
        ]
        
        # Agregar datos de pagos si estÃ¡n disponibles
        if 'pagos' in self.data and self.data['pagos']['total_pagos'] > 0:
            pagos = self.data['pagos']
            data_rows.extend([
                ['Clientes con Pago', f"{pagos.get('clientes_con_pago', 0):,}", '-', 
                 f"Total: ${pagos.get('monto_total', 0):,.0f}"],
                ['Ticket Promedio Pago', f"${pagos.get('ticket_promedio', 0):.2f}", '-', 
                 f"Rango: ${pagos.get('monto_min', 0):.2f} - ${pagos.get('monto_max', 0):,.0f}"]
            ])
        
        # Escribir datos
        for i, row_data in enumerate(data_rows, start=5):
            for j, value in enumerate(row_data, start=1):
                cell = ws.cell(row=i, column=j, value=value)
                if i == 6:  # Encabezados
                    cell.font = Font(bold=True, color=self.COLORS['white'])
                    cell.fill = PatternFill(start_color=self.COLORS['telefonica_light_blue'], 
                                          end_color=self.COLORS['telefonica_light_blue'], fill_type="solid")
        
        # Ajustar anchos
        ws.column_dimensions['A'].width = 25
        ws.column_dimensions['B'].width = 20
        ws.column_dimensions['C'].width = 15
        ws.column_dimensions['D'].width = 40
    
    def _create_excel_analisis_canales(self, wb: openpyxl.Workbook) -> None:
        """Crear hoja de anÃ¡lisis por canales"""
        ws = wb.create_sheet("AnÃ¡lisis por Canal")
        
        # TÃ­tulo
        ws['A1'] = "ANÃLISIS DETALLADO POR CANAL"
        ws['A1'].font = Font(bold=True, size=14, color=self.COLORS['white'])
        ws['A1'].fill = PatternFill(start_color=self.COLORS['telefonica_blue'], 
                                   end_color=self.COLORS['telefonica_blue'], fill_type="solid")
        ws.merge_cells('A1:C1')
        
        # Canal CALL
        call_data = self.data['canal_call']
        call_rows = [
            ['', '', ''],
            ['CANAL CALL', 'VALOR', 'PORCENTAJE'],
            ['Gestiones Totales', f"{call_data.get('total_gestiones', 0):,}", '-'],
            ['Contactos Efectivos', f"{call_data.get('contactos_efectivos', 0):,}", 
             f"{call_data.get('tasa_contactabilidad', 0)}%"],
            ['Contactos No Efectivos', f"{call_data.get('contactos_no_efectivos', 0):,}", '-'],
            ['No Contactos', f"{call_data.get('no_contactos', 0):,}", '-'],
            ['Compromisos', f"{call_data.get('compromisos', 0):,}", 
             f"{call_data.get('tasa_compromiso', 0)}%"],
            ['Monto Compromisos', f"${call_data.get('monto_compromisos', 0):,.0f}", '-'],
            ['DuraciÃ³n Promedio', f"{call_data.get('duracion_promedio', 0):.1f} seg", '-']
        ]
        
        # Escribir datos CALL
        for i, row_data in enumerate(call_rows, start=3):
            for j, value in enumerate(row_data, start=1):
                cell = ws.cell(row=i, column=j, value=value)
                if i == 4:  # Encabezado
                    cell.font = Font(bold=True, color=self.COLORS['white'])
                    cell.fill = PatternFill(start_color=self.COLORS['telefonica_green'], 
                                          end_color=self.COLORS['telefonica_green'], fill_type="solid")
        
        # Canal VOICEBOT
        voicebot_data = self.data['canal_voicebot']
        voicebot_rows = [
            ['', '', ''],
            ['CANAL VOICEBOT', 'VALOR', 'PORCENTAJE'],
            ['Gestiones Totales', f"{voicebot_data.get('total_gestiones', 0):,}", '-'],
            ['Contactos Efectivos', f"{voicebot_data.get('contactos_efectivos', 0):,}", 
             f"{voicebot_data.get('tasa_contactabilidad', 0)}%"],
            ['Compromisos', f"{voicebot_data.get('compromisos', 0):,}", 
             f"{voicebot_data.get('tasa_compromiso', 0)}%"],
        ]
        
        # Escribir datos VOICEBOT
        start_row = len(call_rows) + 5
        for i, row_data in enumerate(voicebot_rows, start=start_row):
            for j, value in enumerate(row_data, start=1):
                cell = ws.cell(row=i, column=j, value=value)
                if i == start_row + 1:  # Encabezado
                    cell.font = Font(bold=True, color=self.COLORS['white'])
                    cell.fill = PatternFill(start_color=self.COLORS['telefonica_orange'], 
                                          end_color=self.COLORS['telefonica_orange'], fill_type="solid")
        
        # Ajustar anchos
        ws.column_dimensions['A'].width = 25
        ws.column_dimensions['B'].width = 20
        ws.column_dimensions['C'].width = 15
    
    def _create_excel_evolucion_diaria(self, wb: openpyxl.Workbook) -> None:
        """Crear hoja de evoluciÃ³n diaria"""
        ws = wb.create_sheet("EvoluciÃ³n Diaria")
        
        # TÃ­tulo
        ws['A1'] = "EVOLUCIÃ“N DIARIA - CONTACTOS EFECTIVOS"
        ws['A1'].font = Font(bold=True, size=14, color=self.COLORS['white'])
        ws['A1'].fill = PatternFill(start_color=self.COLORS['telefonica_blue'], 
                                   end_color=self.COLORS['telefonica_blue'], fill_type="solid")
        ws.merge_cells('A1:H1')
        
        # Encabezados
        headers = ['Fecha', 'CALL Gestiones', 'CALL Contactos', 'VOICEBOT Gestiones', 
                  'VOICEBOT Contactos', 'Total Gestiones', 'Total Contactos', 'Tasa Contactabilidad']
        
        for j, header in enumerate(headers, start=1):
            cell = ws.cell(row=3, column=j, value=header)
            cell.font = Font(bold=True, color=self.COLORS['white'])
            cell.fill = PatternFill(start_color=self.COLORS['telefonica_light_blue'], 
                                   end_color=self.COLORS['telefonica_light_blue'], fill_type="solid")
        
        # Datos diarios
        for i, dia in enumerate(self.data['evolucion_diaria'], start=4):
            row_data = [
                dia['fecha'],
                dia['call_gestiones'],
                dia['call_contactos'],
                dia['voicebot_gestiones'],
                dia['voicebot_contactos'],
                dia['total_gestiones'],
                dia['total_contactos'],
                f"{dia['tasa_contactabilidad']}%"
            ]
            
            for j, value in enumerate(row_data, start=1):
                ws.cell(row=i, column=j, value=value)
        
        # Ajustar anchos
        for col in ['A', 'B', 'C', 'D', 'E', 'F', 'G', 'H']:
            ws.column_dimensions[col].width = 15
    
    def _create_excel_carteras_activas(self, wb: openpyxl.Workbook) -> None:
        """Crear hoja de carteras activas"""
        ws = wb.create_sheet("Carteras Activas")
        
        # TÃ­tulo
        ws['A1'] = "CARTERAS ACTIVAS - PERÃODO ANALIZADO"
        ws['A1'].font = Font(bold=True, size=14, color=self.COLORS['white'])
        ws['A1'].fill = PatternFill(start_color=self.COLORS['telefonica_blue'], 
                                   end_color=self.COLORS['telefonica_blue'], fill_type="solid")
        ws.merge_cells('A1:H1')
        
        # Encabezados
        headers = ['Archivo', 'Tipo Cartera', 'Fecha AsignaciÃ³n', 'Fecha Cierre', 
                  'Clientes Asignados', 'Cuentas', 'DÃ­as Vigencia', 'Estado']
        
        for j, header in enumerate(headers, start=1):
            cell = ws.cell(row=3, column=j, value=header)
            cell.font = Font(bold=True, color=self.COLORS['white'])
            cell.fill = PatternFill(start_color=self.COLORS['telefonica_light_blue'], 
                                   end_color=self.COLORS['telefonica_light_blue'], fill_type="solid")
        
        # Datos de carteras
        for i, cartera in enumerate(self.data['carteras_activas'], start=4):
            row_data = [
                cartera['archivo'],
                cartera['tipo_cartera'],
                cartera['fecha_asignacion'],
                cartera['fecha_cierre'],
                cartera.get('clientes_asignados', 0),
                cartera.get('cuentas_asignadas', 0),
                cartera['dias_vigencia'],
                cartera['estado']
            ]
            
            for j, value in enumerate(row_data, start=1):
                cell = ws.cell(row=i, column=j, value=value)
                # Colorear estado
                if j == 8:  # Columna Estado
                    if value == 'ACTIVA':
                        cell.fill = PatternFill(start_color=self.COLORS['telefonica_green'], 
                                              end_color=self.COLORS['telefonica_green'], fill_type="solid")
                        cell.font = Font(color=self.COLORS['white'])
        
        # Ajustar anchos
        ws.column_dimensions['A'].width = 30
        for col in ['B', 'C', 'D', 'E', 'F', 'G', 'H']:
            ws.column_dimensions[col].width = 15
    
    def _create_excel_kpis_campanias(self, wb: openpyxl.Workbook) -> None:
        """Crear hoja de KPIs por campaÃ±a"""
        ws = wb.create_sheet("KPIs por CampaÃ±a")
        
        # TÃ­tulo
        ws['A1'] = "KPIS DETALLADOS POR CAMPAÃ‘A"
        ws['A1'].font = Font(bold=True, size=14, color=self.COLORS['white'])
        ws['A1'].fill = PatternFill(start_color=self.COLORS['telefonica_blue'], 
                                   end_color=self.COLORS['telefonica_blue'], fill_type="solid")
        ws.merge_cells('A1:H1')
        
        # Encabezados
        headers = ['Archivo', 'Total Gestiones', 'Clientes Gestionados', 'Contactos Efectivos', 
                  'PDPs', 'Monto Compromisos', 'Tasa Contactabilidad', 'Tasa PDP']
        
        for j, header in enumerate(headers, start=1):
            cell = ws.cell(row=3, column=j, value=header)
            cell.font = Font(bold=True, color=self.COLORS['white'])
            cell.fill = PatternFill(start_color=self.COLORS['telefonica_light_blue'], 
                                   end_color=self.COLORS['telefonica_light_blue'], fill_type="solid")
        
        # Datos de KPIs
        for i, kpi in enumerate(self.data['kpis_por_campania'], start=4):
            row_data = [
                kpi.get('archivo', ''),
                kpi.get('total_gestiones', 0),
                kpi.get('clientes_gestionados', 0),
                kpi.get('contactos_efectivos', 0),
                kpi.get('pdps', 0),
                f"${kpi.get('monto_compromisos', 0):,.0f}",
                f"{kpi.get('tasa_contactabilidad', 0)}%",
                f"{kpi.get('tasa_pdp', 0)}%"
            ]
            
            for j, value in enumerate(row_data, start=1):
                ws.cell(row=i, column=j, value=value)
        
        # Ajustar anchos
        ws.column_dimensions['A'].width = 30
        for col in ['B', 'C', 'D', 'E', 'F', 'G', 'H']:
            ws.column_dimensions[col].width = 15
    
    def _create_excel_recomendaciones(self, wb: openpyxl.Workbook) -> None:
        """Crear hoja de recomendaciones"""
        ws = wb.create_sheet("Recomendaciones")
        
        # TÃ­tulo
        ws['A1'] = "RECOMENDACIONES ESTRATÃ‰GICAS"
        ws['A1'].font = Font(bold=True, size=14, color=self.COLORS['white'])
        ws['A1'].fill = PatternFill(start_color=self.COLORS['telefonica_blue'], 
                                   end_color=self.COLORS['telefonica_blue'], fill_type="solid")
        ws.merge_cells('A1:D1')
        
        # Encabezados
        headers = ['CategorÃ­a', 'Prioridad', 'DescripciÃ³n', 'AcciÃ³n Recomendada']
        
        for j, header in enumerate(headers, start=1):
            cell = ws.cell(row=3, column=j, value=header)
            cell.font = Font(bold=True, color=self.COLORS['white'])
            cell.fill = PatternFill(start_color=self.COLORS['telefonica_light_blue'], 
                                   end_color=self.COLORS['telefonica_light_blue'], fill_type="solid")
        
        # Datos de recomendaciones
        for i, rec in enumerate(self.data['recomendaciones'], start=4):
            row_data = [
                rec.get('categoria', ''),
                rec.get('prioridad', ''),
                rec.get('descripcion', ''),
                rec.get('accion', '')
            ]
            
            for j, value in enumerate(row_data, start=1):
                cell = ws.cell(row=i, column=j, value=value)
                # Colorear por prioridad
                if j == 2:  # Columna Prioridad
                    if value == 'Alta':
                        cell.fill = PatternFill(start_color='FF6B6B', end_color='FF6B6B', fill_type="solid")
                        cell.font = Font(color=self.COLORS['white'])
                    elif value == 'Media':
                        cell.fill = PatternFill(start_color='FFE66D', end_color='FFE66D', fill_type="solid")
        
        # Ajustar anchos
        ws.column_dimensions['A'].width = 20
        ws.column_dimensions['B'].width = 12
        ws.column_dimensions['C'].width = 40
        ws.column_dimensions['D'].width = 50
    
    def generate_powerpoint_report(self, output_path: str = None) -> str:
        """
        Generar presentaciÃ³n PowerPoint ejecutiva
        """
        if output_path is None:
            timestamp = self.fecha_generacion.strftime('%Y%m%d_%H%M%S')
            output_path = f"Presentacion_Semanal_Telefonica_{timestamp}.pptx"
        
        logger.info(f"ðŸ“ˆ Generando presentaciÃ³n PowerPoint: {output_path}")
        
        try:
            # Crear presentaciÃ³n
            prs = Presentation()
            
            # Generar slides
            self._create_ppt_portada(prs)
            self._create_ppt_resumen_ejecutivo(prs)
            self._create_ppt_analisis_canales(prs)
            self._create_ppt_evolucion_temporal(prs)
            self._create_ppt_carteras_activas(prs)
            self._create_ppt_recomendaciones(prs)
            
            # Guardar presentaciÃ³n
            prs.save(output_path)
            logger.info(f"âœ… PresentaciÃ³n PowerPoint generada exitosamente: {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"âŒ Error generando PowerPoint: {str(e)}")
            raise
    
    def _create_ppt_portada(self, prs: Presentation) -> None:
        """Crear slide de portada"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "INFORME SEMANAL DE GESTIÃ“N"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = f"TelefÃ³nica del PerÃº\\n{self.periodo_str}\\nSistema de Cobranza Automatizado"
        subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    def _create_ppt_resumen_ejecutivo(self, prs: Presentation) -> None:
        """Crear slide de resumen ejecutivo"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "RESUMEN EJECUTIVO"
        
        resumen = self.data['resumen_ejecutivo']
        tf = content.text_frame
        tf.text = f"â€¢ {resumen.get('total_gestiones', 0):,} gestiones totales realizadas"
        
        # Agregar pÃ¡rrafos adicionales
        paragraphs_data = [
            f"â€¢ {resumen.get('total_contactos_efectivos', 0):,} contactos efectivos ({resumen.get('tasa_contactabilidad_global', 0)}%)",
            f"â€¢ {resumen.get('total_compromisos', 0):,} compromisos obtenidos ({resumen.get('tasa_compromiso_global', 0)}%)",
            f"â€¢ ${resumen.get('monto_compromisos_call', 0):,.0f} en compromisos CALL",
            f"â€¢ {resumen.get('clientes_unicos_total', 0):,} clientes Ãºnicos gestionados"
        ]
        
        # Agregar informaciÃ³n de pagos si disponible
        if 'pagos' in self.data and self.data['pagos']['total_pagos'] > 0:
            pagos = self.data['pagos']
            paragraphs_data.append(f"â€¢ ${pagos.get('monto_total', 0):,.0f} en pagos procesados")
        
        for para_text in paragraphs_data:
            p = tf.add_paragraph()
            p.text = para_text
            p.font.size = Pt(16)
    
    def _create_ppt_analisis_canales(self, prs: Presentation) -> None:
        """Crear slide de anÃ¡lisis por canales"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # TÃ­tulo
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        tf.text = "ANÃLISIS POR CANAL"
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        
        # Canal CALL
        call_data = self.data['canal_call']
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(4)
        height = Inches(3.5)
        call_box = slide.shapes.add_textbox(left, top, width, height)
        tf_call = call_box.text_frame
        tf_call.text = "ðŸ“ž CANAL CALL"
        tf_call.paragraphs[0].font.size = Pt(20)
        tf_call.paragraphs[0].font.bold = True
        
        call_bullets = [
            f"â€¢ {call_data.get('total_gestiones', 0):,} gestiones",
            f"â€¢ {call_data.get('tasa_contactabilidad', 0)}% contactabilidad",
            f"â€¢ {call_data.get('compromisos', 0):,} compromisos",
            f"â€¢ ${call_data.get('monto_compromisos', 0):,.0f} monto"
        ]
        
        for bullet in call_bullets:
            p = tf_call.add_paragraph()
            p.text = bullet
            p.font.size = Pt(14)
        
        # Canal VOICEBOT
        voicebot_data = self.data['canal_voicebot']
        left = Inches(5)
        top = Inches(2)
        width = Inches(4)
        height = Inches(3.5)
        vb_box = slide.shapes.add_textbox(left, top, width, height)
        tf_vb = vb_box.text_frame
        tf_vb.text = "ðŸ¤– CANAL VOICEBOT"
        tf_vb.paragraphs[0].font.size = Pt(20)
        tf_vb.paragraphs[0].font.bold = True
        
        vb_bullets = [
            f"â€¢ {voicebot_data.get('total_gestiones', 0):,} gestiones",
            f"â€¢ {voicebot_data.get('tasa_contactabilidad', 0)}% contactabilidad",
            f"â€¢ {voicebot_data.get('compromisos', 0):,} compromisos"
        ]
        
        for bullet in vb_bullets:
            p = tf_vb.add_paragraph()
            p.text = bullet
            p.font.size = Pt(14)
    
    def _create_ppt_evolucion_temporal(self, prs: Presentation) -> None:
        """Crear slide de evoluciÃ³n temporal"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "EVOLUCIÃ“N TEMPORAL"
        
        tf = content.text_frame
        tf.text = "Tendencias de Contactabilidad por DÃ­a:"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        
        # Encontrar mejor y peor dÃ­a
        if self.data['evolucion_diaria']:
            mejor_dia = max(self.data['evolucion_diaria'], key=lambda x: x['tasa_contactabilidad'])
            peor_dia = min(self.data['evolucion_diaria'], key=lambda x: x['tasa_contactabilidad'])
            
            insights = [
                f"â€¢ Mejor dÃ­a: {mejor_dia['fecha']} ({mejor_dia['tasa_contactabilidad']}% contactabilidad)",
                f"â€¢ Menor dÃ­a: {peor_dia['fecha']} ({peor_dia['tasa_contactabilidad']}% contactabilidad)",
                f"â€¢ Total dÃ­as analizados: {len(self.data['evolucion_diaria'])}",
                f"â€¢ Promedio contactabilidad: {sum(d['tasa_contactabilidad'] for d in self.data['evolucion_diaria']) / len(self.data['evolucion_diaria']):.1f}%"
            ]
            
            for insight in insights:
                p = tf.add_paragraph()
                p.text = insight
                p.font.size = Pt(14)
    
    def _create_ppt_carteras_activas(self, prs: Presentation) -> None:
        """Crear slide de carteras activas"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "CARTERAS ACTIVAS"
        
        tf = content.text_frame
        total_clientes = sum(c.get('clientes_asignados', 0) for c in self.data['carteras_activas'])
        tf.text = f"â€¢ {total_clientes:,} clientes asignados total"
        tf.paragraphs[0].font.size = Pt(16)
        
        # Resumen por tipo de cartera
        cartera_summary = {}
        for cartera in self.data['carteras_activas']:
            tipo = cartera['tipo_cartera']
            clientes = cartera.get('clientes_asignados', 0)
            if tipo in cartera_summary:
                cartera_summary[tipo] += clientes
            else:
                cartera_summary[tipo] = clientes
        
        for tipo, clientes in cartera_summary.items():
            p = tf.add_paragraph()
            p.text = f"â€¢ {tipo}: {clientes:,} clientes"
            p.font.size = Pt(14)
        
        # InformaciÃ³n de vigencias
        activas = len([c for c in self.data['carteras_activas'] if c['estado'] == 'ACTIVA'])
        p = tf.add_paragraph()
        p.text = f"â€¢ {activas} carteras activas de {len(self.data['carteras_activas'])} total"
        p.font.size = Pt(14)
    
    def _create_ppt_recomendaciones(self, prs: Presentation) -> None:
        """Crear slide de recomendaciones"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "RECOMENDACIONES ESTRATÃ‰GICAS"
        
        tf = content.text_frame
        
        if not self.data['recomendaciones']:
            tf.text = "â€¢ Mantener monitoreo continuo de KPIs"
            tf.paragraphs[0].font.size = Pt(16)
            
            general_recs = [
                "â€¢ Optimizar distribuciÃ³n de cartera entre canales",
                "â€¢ Implementar mejoras continuas en procesos",
                "â€¢ Desarrollar anÃ¡lisis predictivos"
            ]
            
            for rec in general_recs:
                p = tf.add_paragraph()
                p.text = rec
                p.font.size = Pt(14)
        else:
            # Tomar las 5 recomendaciones mÃ¡s importantes
            top_recommendations = self.data['recomendaciones'][:5]
            
            for i, rec in enumerate(top_recommendations):
                if i == 0:
                    tf.text = f"â€¢ {rec['descripcion']}"
                    tf.paragraphs[0].font.size = Pt(14)
                else:
                    p = tf.add_paragraph()
                    p.text = f"â€¢ {rec['descripcion']}"
                    p.font.size = Pt(14)
    
    def generate_complete_report(self, output_dir: str = None) -> Tuple[str, str]:
        """
        Generar reporte completo (Excel + PowerPoint)
        """
        if output_dir is None:
            output_dir = tempfile.gettempdir()
        
        timestamp = self.fecha_generacion.strftime('%Y%m%d_%H%M%S')
        
        excel_path = os.path.join(output_dir, f"Informe_Semanal_Telefonica_{timestamp}.xlsx")
        ppt_path = os.path.join(output_dir, f"Presentacion_Semanal_Telefonica_{timestamp}.pptx")
        
        # Generar ambos reportes
        excel_file = self.generate_excel_report(excel_path)
        ppt_file = self.generate_powerpoint_report(ppt_path)
        
        logger.info(f"ðŸŽ‰ Reportes completos generados en: {output_dir}")
        
        return excel_file, ppt_file
